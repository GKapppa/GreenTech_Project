from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1B, 0x5E, 0x20)
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(1))
        tf = sub_box.text_frame
        tf.paragraphs[0].text = subtitle
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xFF, 0xCC)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, image_path=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1B, 0x5E, 0x20)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if image_path:
        try:
            slide.shapes.add_picture(image_path, Inches(7.5), Inches(1.5), width=Inches(5.3))
        except:
            pass

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.8), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(12)

    return slide

def add_two_column_slide(prs, title, left_content, right_content, left_title="", right_title=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1B, 0x5E, 0x20)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if left_title:
        lt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.5))
        tf = lt_box.text_frame
        tf.paragraphs[0].text = left_title
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.8), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(8)

    if right_title:
        rt_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.5))
        tf = rt_box.text_frame
        tf.paragraphs[0].text = right_title
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)

    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(5.8), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(8)

    return slide

add_title_slide(
    prs,
    "Mentor IA GreenTech",
    "Sistema de Induccion Tecnica con Agentes Inteligentes\n\nAsignatura: Ingenieria de Soluciones con IA - ISY0101\nDuoc UC"
)

add_content_slide(prs, "1. Introduccion y Caso de Negocio", [
    "PROBLEMA: Los nuevos tecnicos necesitan consultar manuales tecnicos extensos sobre sistemas fotovoltaicos",
    "Este proceso toma entre 30-60 minutos por tema",
    "No siempre es facil encontrar la informacion correcta",
    "",
    "SOLUCION: Asistente virtual inteligente que permite consultar",
    "los manuales mediante lenguaje natural",
    "",
    "BENEFICIOS:",
    "  - Reduce tiempo de capacitacion de horas a segundos",
    "  - Respuestas con rigor tecnico basado en manuales oficiales",
    "  - Disponible 24/7 para consulta de tecnicos",
    "  - Trazabilidad completa de todas las consultas"
])

add_content_slide(prs, "2. Arquitectura del Sistema", [
    "ARQUITECTURA DE 4 CAPAS:",
    "",
    "Capa 1 - Interfaz de Usuario:",
    "  Streamlit (aplicacion web con chat intuitivo)",
    "",
    "Capa 2 - Agente Orquestador:",
    "  Coordina decisiones, clasifica intenciones",
    "",
    "Capa 3 - Motor RAG:",
    "  MongoDB Atlas (busqueda vectorial semantica)",
    "",
    "Capa 4 - Modelo de Lenguaje:",
    "  Groq con Llama 3.3 (70B parametros)"
])

add_content_slide(prs, "3. Diagrama de Arquitectura", [
    "FLUJO DE DATOS:",
    "",
    "Usuario -> Streamlit Chat -> Agente Orquestador",
    "                                |",
    "                                v",
    "                    Planner (clasifica intencion)",
    "                                |",
    "                +---------------+---------------+",
    "                |               |               |",
    "                v               v               v",
    "          Documents      Memory         Report",
    "          (RAG)          (JSON)       Generator",
    "                |               |               |",
    "                +---------------+---------------+",
    "                                |",
    "                                v",
    "                    Groq LLM (genera respuesta)",
    "                                |",
    "                                v",
    "                    Respuesta al Usuario"
])

add_content_slide(prs, "4. Agente Orquestador - Planner", [
    "El Planner clasifica la pregunta en 5 intenciones:",
    "",
    "1. CONSULTA_SIMPLE: Responde sin buscar documentos",
    "   (saludos, preguntas generales)",
    "",
    "2. CONSULTA_TECNICA: Busca en manuales vectorizados",
    "   (terminos tecnicos del dominio)",
    "",
    "3. SOLICITUD_REPORTE: Genera reporte estructurado",
    "   (cuando el usuario pide 'informe' o 'reporte')",
    "",
    "4. CONTINUIDAD_CONTEXTO: Usa memoria de conversacion",
    "   (referencias a conversaciones anteriores)",
    "",
    "5. INFORMACION_INSUFICIENTE: Pide mas detalles",
    "   (preguntas muy breves o ambiguas)"
])

add_content_slide(prs, "5. Motor RAG - Recuperacion Semantica", [
    "COMO FUNCIONA LA RECUPERACION:",
    "",
    "1. Documentos PDF se fragmentan en chunks (~500 caracteres)",
    "",
    "2. Cada chunk se convierte en embedding numerico",
    "   usando all-MiniLM-L6-v2 (384 dimensiones)",
    "",
    "3. Los vectores se almacenan en MongoDB Atlas",
    "   en coleccion 'manuals_vectors'",
    "",
    "4. Cuando el usuario pregunta:",
    "   - La pregunta se Embedda",
    "   - Se busca en MongoDB los chunks mas similares",
    "   - Similitud coseno determina relevancia",
    "",
    "5. Contexto recuperado se envia al LLM",
    "   para generar respuesta final"
])

add_content_slide(prs, "6. Sistema de Memoria", [
    "DOS TIPOS DE MEMORIA:",
    "",
    "MEMORIA CORTA (Session State):",
    "  - Mantiene historial mientras la app esta abierta",
    "  - streamlit session_state",
    "  - Acceso rapido a ultimos mensajes",
    "",
    "MEMORIA LARGA (JSON):",
    "  - Guarda resumenes de conversaciones anteriores",
    "  - Archivo JSON local en data/memory/",
    "  - Permite continuidad en sesiones futuras",
    "",
    "BENEFICIO: El agente recuerda contexto previo",
    "y puede responder preguntas de seguimiento"
])

add_two_column_slide(
    prs,
    "7. Observabilidad y Monitoreo",
    [
        "METRICAS MONITOREADAS:",
        "",
        "1. Latencias por fase:",
        "   - Planner, Retrieval, Generation",
        "",
        "2. Tokens consumidos:",
        "   - Input, output, costo USD",
        "",
        "3. Precision RAG:",
        "   - Similitud coseno pregunta vs contexto",
        "",
        "4. Calidad de respuestas:",
        "   - Evalua si tiene contenido util",
        "",
        "5. Metricas de sistema:",
        "   - CPU, memoria RAM (psutil)",
    ],
    [
        "DASHBOARD EN STREAMLIT:",
        "",
        "- Tarjetas KPI en tiempo real",
        "- Graficos de tendencia",
        "- Percentiles P50, P95, P99",
        "- Deteccion automatica de anomalias",
        "- Patron de errores por intencion",
        "- Tasa de error en ventanas moviles",
        "- Explorador de logs con filtros",
        "- Exportacion a CSV",
        "",
        "ANOMALY DETECTOR:",
        "- Detecta latencias > P95",
        "- Identifica rafagas de errores",
        "- Genera recomendaciones",
    ],
    "Metricas", "Dashboard"
)

add_content_slide(prs, "8. Deteccion de Anomalias", [
    "COMO DETECTAMOS ANOMALIAS:",
    "",
    "METODO DEL PERCENTIL 95:",
    "  - Se calculan los percentiles de latencia",
    "  - Si una latencia supera P95, se marca como anomalia",
    "",
    "RAFAGAS DE ERRORES:",
    "  - Se detectan cuando 4 de 5 consultas consecutivas fallan",
    "  - Indica problema sistematico",
    "",
    "PATRONES DE ERROR POR INTENCION:",
    "  - Agrupa errores por tipo de consulta",
    "  - Identifica cual intencion genera mas problemas",
    "",
    "RECOMENDACIONES AUTOMATICAS:",
    "  - El sistema genera sugerencias segun los patrones",
    "  - Ej: 'Latencia promedio alta, considerar caching'"
])

add_two_column_slide(
    prs,
    "9. Seguridad y Compliance",
    [
        "TRES CAPAS DE SEGURIDAD:",
        "",
        "1. Enmascaramiento PII:",
        "   - Emails, telefonos, credenciales",
        "   - Regex para deteccion automatica",
        "",
        "2. Deteccion Inyeccion Prompt:",
        "   - Lista de keywords sospechosas",
        "   - 'ignore instructions', 'jailbreak'",
        "   - Bloqueo inmediato sin consumir LLM",
        "",
        "3. Auditoria Inmutable:",
        "   - Logs con timestamp UTC",
        "   - Session ID para tracking",
    ],
    [
        "NORMATIVA CUMPLIDA:",
        "",
        "- Ley 19.628 Chile (Datos Personales)",
        "- GDPR / LGPD (principios aplicados)",
        "- NIST AI Risk Framework",
        "",
        "POLITICA DE RETENCION:",
        "- Logs: maximo 90 dias",
        "- Backups cifrados AES-256",
        "- Consentimiento registrado",
    ],
    "Capas de Proteccion", "Compliance"
)

add_content_slide(prs, "10. Caching Semantico", [
    "OPTIMIZACION CON CACHE:",
    "",
    "PROBLEMA: Consultas similares generan costos y latencia innecesarios",
    "",
    "SOLUCION: SemanticCache con TF-IDF",
    "  - Almacena respuestas anteriores",
    "  - Usa similitud coseno para detectar consultas similares",
    "  - Umbral configurable (default 0.92)",
    "",
    "BENEFICIOS MEDIDOS:",
    "  - Latencia reducida a < 0.05s para cache hits",
    "  - Tokens ahorrados en consultas repetitivas",
    "  - Costo operacional reducido",
    "",
    "EVICCION AUTOMATICA:",
    "  - LRU (Least Recently Used)",
    "  - Maximo 500 entradas configurables"
])

add_content_slide(prs, "11. Tecnologias Utilizadas", [
    "BACKEND Y AGENTE:",
    "  - Python 3.14",
    "  - LangChain (framework de agentes)",
    "  - LangChain Groq (integracion LLM)",
    "",
    "VECTOR STORE:",
    "  - MongoDB Atlas (base de datos)",
    "  - HuggingFace Embeddings (all-MiniLM-L6-v2)",
    "",
    "INTERFAZ:",
    "  - Streamlit (framework web Python)",
    "  - Pandas (manipulacion de datos)",
    "",
    "MONITOREO:",
    "  - psutil (metricas de sistema)",
    "  - scikit-learn (TF-IDF para cache)",
    "",
    "LLM:",
    "  - Groq API (Llama 3.3 70B)",
])

add_content_slide(prs, "12. Limitaciones y Mejoras Futuras", [
    "LIMITACIONES ACTUALES:",
    "",
    "  1. Precision RAG depende de calidad de embeddings",
    "  2. Clasificacion por reglas, no por ML",
    "  3. Sin pruebas automatizadas E2E",
    "  4. Memoria larga simple (sin compresion avanzada)",
    "",
    "MEJORAS PROPUESTAS:",
    "",
    "  1. Caching semantico (YA IMPLEMENTADO)",
    "     - Reduce latencia y costo",
    "",
    "  2. Modelo SLM para tareas auxiliares",
    "     - Llama-3-8B para Planner y Guardrails",
    "     - Reduce latencia inicial",
    "",
    "  3. Busqueda hibrida con reranking",
    "     - Vectorial + BM25 + Cohere Rerank",
    "     - Precision > 95%"
])

add_content_slide(prs, "13. Caso de Negocio - ROI", [
    "ANALISIS COSTO-BENEFICIO:",
    "",
    "COSTO POR CONSULTA:",
    "  - Promedio: $0.0005 USD",
    "  - Para 1000 consultas/dia: $0.50 USD/dia",
    "  - Proyeccion mensual: $15 USD/mes",
    "",
    "AHORRO EN CAPACITACION:",
    "  - Tiempo manual: 30-60 min por tema",
    "  - Con Mentor IA: segundos",
    "  - Tecnicos nuevos: ~20 al mes",
    "  - Ahorro estimado: 10-20 horas/mes",
    "",
    "REDUCCION DE ERRORES:",
    "  - Consultas precisas basadas en manuales",
    "  - Menos interpretaciones incorrectas",
    "  - Mayor seguridad en instalaciones"
])

add_content_slide(prs, "14. Demo - Sistema en Funcionamiento", [
    "PANTALLA PRINCIPAL:",
    "",
    "  - Menu lateral con modulos de aprendizaje",
    "  - Preguntas sugeridas clickeables",
    "  - Area de chat central",
    "",
    "PESTAÑAS:",
    "  1. Tutoria Interactiva (Chat)",
    "  2. Dashboard de Observabilidad",
    "  3. Auditoria y Compliance",
    "",
    "DURANTE LA DEMO MOSTRAREMOS:",
    "  - Pregunta tecnica: 'Explicame el efecto fotoelectrico'",
    "  - Advertencia automatica de seguridad",
    "  - Dashboard con metricas en tiempo real",
    "  - Deteccion de anomalias y recomendaciones"
])

add_title_slide(
    prs,
    "Gracias por su atencion",
    "Preguntas?\n\nProyecto Mentor IA GreenTech\nISY0101 - Duoc UC"
)

prs.save("Presentacion_MentorIA_GreenTech.pptx")
print("Presentacion guardada: Presentacion_MentorIA_GreenTech.pptx")
